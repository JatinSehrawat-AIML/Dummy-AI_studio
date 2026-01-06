# loaders/ppt_loader.py
from pptx import Presentation
from pathlib import Path


def load_ppt(path: str) -> str:
    ppt_path = Path(path)

    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT file not found at: {ppt_path}")

    prs = Presentation(ppt_path)
    slides_text = []

    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)

        if slide_content:
            slides_text.append(" ".join(slide_content))

    text = "\n".join(slides_text)

    if not text.strip():
        raise ValueError("No readable text found in PPT")

    return text
